# analyzer.py
# Read a pitch-deck PDF/PPTX → Gemini returns strict JSON → score → INVEST/DEFER/PASS

import os, re, json, io, mimetypes
from typing import Any, Dict, Optional
from google import genai
from google.genai import types

# --- dev key (inline as you requested). Env var overrides if set. ---
# API_KEY = os.getenv("GEMINI_API_KEY") or "NONE"
client = genai.Client(api_key="AIzaSyCAnjBDLQqnHtJkc4su1iXwmsjwOzHUtKE")

EXTRACTION_PROMPT = """
You are analyzing a STARTUP PITCH DECK (PDF or PPTX).
Return ONLY valid JSON (no markdown, no fences, no commentary). If a value is missing, set it to null.
Do NOT invent numbers; extract only if explicitly present. Normalize money to plain numbers in USD (e.g., "$1.2M" → 1200000).
Top-level keys MUST appear exactly once.

Schema:
{
  "summary": "",                       // 3–6 sentences: what they do, who for, traction, why now
  "startup": {
    "name": null, "website": null, "sector": null, "subsector": null,
    "stage": null, "hq_country": null, "founded_year": null
  },
  "team": {
    "founders": [], "team_size": null
  },
  "traction": {
    "mrr_usd": null, "arr_usd": null, "growth_mom_pct": null,
    "customers": null, "churn_pct": null, "retention_pct": null
  },
  "unit_economics": {
    "cac_usd": null, "ltv_usd": null, "gross_margin_pct": null,
    "burn_rate_usd_per_month": null, "runway_months": null
  },
  "market": {
    "tam_usd": null, "sam_usd": null, "som_usd": null, "competitors_count": null
  },
  "round": {
    "seeking_usd": null, "pre_money_valuation_usd": null, "existing_investors": []
  },
  "risks": { "top_3_risks": [] }
}
Output exactly one JSON object.
"""

def _num(x: Any) -> Optional[float]:
    if x is None: return None
    if isinstance(x, (int, float)): return float(x)
    try:
        s = str(x).strip()
        s = re.sub(r"[,\s$]", "", s, flags=re.I)
        s = s.lower().replace("usd", "")
        mult = 1.0
        if s.endswith("k"): mult, s = 1_000.0, s[:-1]
        elif s.endswith("m"): mult, s = 1_000_000.0, s[:-1]
        elif s.endswith("b"): mult, s = 1_000_000_000.0, s[:-1]
        return float(s) * mult
    except:
        return None

def score_and_verdict(data: Dict[str, Any]) -> Dict[str, Any]:
    s = 50
    reasons = []

    tr  = data.get("traction") or {}
    ue  = data.get("unit_economics") or {}
    mkt = data.get("market") or {}
    team= data.get("team") or {}

    arr   = _num(tr.get("arr_usd"))
    mrr   = _num(tr.get("mrr_usd"))
    mom   = _num(tr.get("growth_mom_pct"))
    churn = _num(tr.get("churn_pct"))
    ret   = _num(tr.get("retention_pct"))

    cac   = _num(ue.get("cac_usd"))
    ltv   = _num(ue.get("ltv_usd"))
    gm    = _num(ue.get("gross_margin_pct"))
    burn  = _num(ue.get("burn_rate_usd_per_month"))
    runway= _num(ue.get("runway_months"))

    tam   = _num(mkt.get("tam_usd"))
    tsize = _num(team.get("team_size"))
    founders = team.get("founders") or []

    # ARR / MRR
    if arr is not None:
        if arr >= 1_000_000: s += 20; reasons.append(f"Strong ARR ${int(arr):,}.")
        elif arr >= 250_000: s += 10; reasons.append(f"Decent ARR ${int(arr):,}.")
        else: s -= 6; reasons.append("Low ARR.")
    elif mrr is not None:
        if mrr >= 80_000: s += 12; reasons.append(f"Solid MRR ${int(mrr):,}/mo.")
        elif mrr >= 20_000: s += 6; reasons.append("OK MRR.")
        else: s -= 5; reasons.append("Low MRR.")

    # Growth
    if mom is not None:
        if mom >= 15: s += 12; reasons.append(f"Healthy growth {mom}% MoM.")
        elif mom >= 5: s += 6; reasons.append("Moderate growth.")
        elif mom < 0: s -= 12; reasons.append("Negative growth.")
        else: s -= 3; reasons.append("Low growth.")

    # Churn / Retention
    if churn is not None and churn > 6: s -= 8; reasons.append(f"High churn {churn}%.")
    if ret is not None:
        if ret >= 85: s += 4; reasons.append(f"Good retention {ret}%.")
        elif ret < 70: s -= 6; reasons.append(f"Weak retention {ret}%.")

    # LTV/CAC
    if ltv and cac and cac > 0:
        ratio = ltv / cac
        if ratio >= 3: s += 15; reasons.append(f"Efficient LTV/CAC {ratio:.1f}.")
        elif ratio >= 2: s += 7; reasons.append(f"Acceptable LTV/CAC {ratio:.1f}.")
        else: s -= 10; reasons.append(f"Poor LTV/CAC {ratio:.1f}.")

    # Gross margin
    if gm is not None:
        if gm >= 70: s += 4; reasons.append("High gross margin.")
        elif gm < 40: s -= 5; reasons.append("Low gross margin.")

    # Burn / runway
    if burn is not None and runway is not None:
        if runway < 6 and burn > 0: s -= 8; reasons.append("Short runway (<6 months).")
        elif runway >= 12: s += 4; reasons.append("Comfortable runway (≥12 months).")

    # TAM
    if tam is not None and tam >= 1_000_000_000: s += 5; reasons.append("Large TAM (≥$1B).")

    # Team
    if tsize is not None and tsize >= 8: s += 3; reasons.append("Adequate team size.")
    if any(re.search(r"\b(exit|acqui|unicorn|ex-|FAANG|IPO)\b", str(f), re.I) for f in founders):
        s += 6; reasons.append("Founder signals (exits/top-tier).")

    s = max(0, min(100, s))
    verdict = "INVEST" if s >= 65 else ("DEFER" if s >= 50 else "PASS")
    return {"score": s, "verdict": verdict, "reasons": reasons}

def _extract_text_from_pptx(path: str) -> str:
    # Fallback path for PPTX (Gemini can work with text if you don't want to convert to PDF).
    try:
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks).strip()
    except Exception as e:
        raise RuntimeError(f"Failed to read PPTX: {e}")

def analyze_pitchdeck(file_path: str) -> Dict[str, Any]:
    """
    Pass a local file path (.pdf or .pptx). Returns:
        {"extracted": <schema>, "decision": {"score": int, "verdict": str, "reasons": []}}
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(file_path)

    ext = os.path.splitext(file_path)[1].lower()
    mime, _ = mimetypes.guess_type(file_path)
    mime = mime or ("application/pdf" if ext == ".pdf" else None)

    # Build Gemini input
    parts = [EXTRACTION_PROMPT]
    if ext == ".pdf":
        with open(file_path, "rb") as f:
            pdf_bytes = f.read()
        parts.append(types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf"))
    elif ext == ".pptx":
        # Use extracted text (simple + works well for text-heavy decks)
        text = _extract_text_from_pptx(file_path)
        parts.append(text if text else "(empty file)")
    else:
        raise ValueError("Unsupported file type. Use .pdf or .pptx")

    # Call Gemini → enforce JSON
    try:
        resp = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=parts,
            generation_config={"response_mime_type": "application/json", "temperature": 0.1},
        )
        text = resp.text or ""
        data = json.loads(text)
    except TypeError:
        # Older google-genai versions: no generation_config
        resp = client.models.generate_content(model="gemini-2.5-flash", contents=parts)
        text = (resp.text or "").strip()
        text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.I|re.M)   # strip code fences
        text = re.sub(r",(\s*[}\]])", r"\1", text)                             # trailing commas
        if "{" in text and "}" in text:
            text = text[text.find("{"): text.rfind("}")+1]
        data = json.loads(text)

    decision = score_and_verdict(data)
    return {"extracted": data, "decision": decision}

if __name__ == "__main__":
    # quick manual test: set a path, print result
    sample_path = r"C:\Users\Armaan\Desktop\Jay-Mehta-Pitch-Deck.pdf"
    print(json.dumps(analyze_pitchdeck(sample_path), indent=2))
